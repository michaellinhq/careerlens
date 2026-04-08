"""
CareerLens Pitch Deck Generator
Generates Chinese and German business plan presentations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BLUE = RGBColor(37, 99, 235)      # #2563EB
DARK = RGBColor(15, 23, 42)       # #0F172A
SLATE = RGBColor(71, 85, 105)     # #475569
LIGHT = RGBColor(248, 250, 252)   # #F8FAFC
WHITE = RGBColor(255, 255, 255)
EMERALD = RGBColor(5, 150, 105)   # #059669
RED = RGBColor(220, 38, 38)       # #DC2626
AMBER = RGBColor(217, 119, 6)     # #D97706

def add_bg(slide, color=LIGHT):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=14, color=SLATE, bold=False, space_before=6, font_name='Calibri'):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = Pt(space_before)
    return p

def add_rect(slide, left, top, width, height, fill_color):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def build_deck(lang='zh'):
    """Build a complete pitch deck in the specified language."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    is_zh = lang == 'zh'

    # ==================== SLIDE 1: Title ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, BLUE)

    add_textbox(slide, 1, 1.5, 11, 1.2,
        'CareerLens' if not is_zh else 'CareerLens 职业透镜',
        font_size=48, color=WHITE, bold=True)

    add_textbox(slide, 1, 2.8, 11, 0.8,
        'AI驱动的工程师技能转型平台' if is_zh else
        'KI-gesteuerte Skill-Transformationsplattform für Ingenieure',
        font_size=24, color=RGBColor(191, 219, 254))

    tf = add_textbox(slide, 1, 4.2, 11, 2,
        '商业计划书 · 2026' if is_zh else
        'Geschäftsplan · 2026',
        font_size=18, color=RGBColor(191, 219, 254))
    add_para(tf, '', font_size=12)
    add_para(tf,
        'Michael Lin (林海青) · Magna 汽车电子质量管理 10年+' if is_zh else
        'Michael Lin (林海青) · 10+ Jahre Automotive Quality Management bei Magna',
        font_size=14, color=RGBColor(191, 219, 254))
    add_para(tf,
        '联系方式: michaellinhq@gmail.com · github.com/michaellinhq/careerlens' if is_zh else
        'Kontakt: michaellinhq@gmail.com · github.com/michaellinhq/careerlens',
        font_size=12, color=RGBColor(148, 163, 184))

    # ==================== SLIDE 2: Problem ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '问题：全球工程师正在经历前所未有的技能危机' if is_zh else
        'Problem: Ingenieure erleben eine beispiellose Qualifikationskrise',
        font_size=32, color=DARK, bold=True)

    # Left column - Germany
    add_rect(slide, 1, 1.8, 5.4, 4.8, WHITE)
    tf = add_textbox(slide, 1.3, 2, 5, 0.5,
        '🇩🇪 德国：300万+汽车工人面临失业' if is_zh else
        '🇩🇪 Deutschland: 300.000+ Autoindustrie-Arbeiter bedroht',
        font_size=18, color=RED, bold=True)
    items_de = [
        ('VW 裁员 30,000人 · Continental 裁员 7,000人 · ZF 裁员 12,000人' if is_zh else
         'VW: 30.000 · Continental: 7.000 · ZF: 12.000 Stellenabbau'),
        ('传统内燃机技能快速贬值，电动化/数字化人才严重短缺' if is_zh else
         'Traditionelle Verbrenner-Skills verlieren rapide an Wert'),
        ('Arbeitsamt (BA) 每年拨款数十亿欧元用于再培训' if is_zh else
         'BA investiert Milliarden jährlich in Umschulung'),
        ('但工具落后：纸质问卷 + 人工顾问看30秒简历' if is_zh else
         'Aber veraltete Tools: Papierfragebögen + 30-Sekunden-CV-Check'),
    ]
    for item in items_de:
        add_para(tf, f'• {item}', font_size=13, color=SLATE, space_before=10)

    # Right column - China
    add_rect(slide, 6.8, 1.8, 5.4, 4.8, WHITE)
    tf = add_textbox(slide, 7.1, 2, 5, 0.5,
        '🇨🇳 中国：1000万+中年工程师就业困难' if is_zh else
        '🇨🇳 China: 10 Mio.+ erfahrene Ingenieure in Karrierekrise',
        font_size=18, color=RED, bold=True)
    items_cn = [
        ('制造业PMI持续收缩，35岁危机加剧' if is_zh else
         'Industrie-PMI schrumpft, „35-Jahre-Krise" verschärft sich'),
        ('2.8亿自由职业者缺乏转型指导' if is_zh else
         '280 Mio. Freelancer ohne Transformationsberatung'),
        ('AI替代风险：传统质量/测试岗位首当其冲' if is_zh else
         'KI-Risiko: Traditionelle QA/Test-Jobs zuerst betroffen'),
        ('现有平台(BOSS/LinkedIn)只做招聘，不做转型' if is_zh else
         'Bestehende Plattformen (LinkedIn/StepStone) vermitteln nur, transformieren nicht'),
    ]
    for item in items_cn:
        add_para(tf, f'• {item}', font_size=13, color=SLATE, space_before=10)

    # Bottom insight
    add_textbox(slide, 1, 6.8, 11, 0.5,
        '核心矛盾：数百万工程师有可迁移技能，但不知道自己"离新身份只差1-3个技能"' if is_zh else
        'Kernproblem: Millionen Ingenieure haben transferierbare Skills, wissen aber nicht, dass sie nur 1-3 Skills von einer neuen Karriere entfernt sind',
        font_size=14, color=BLUE, bold=True)

    # ==================== SLIDE 3: Solution ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '解决方案：AI驱动的"人才错配修复引擎"' if is_zh else
        'Lösung: KI-gesteuerte „Talent-Mismatch-Repair-Engine"',
        font_size=32, color=DARK, bold=True)

    add_textbox(slide, 1, 1.5, 11, 0.6,
        'CareerLens 不是又一个招聘网站。我们修复"技能-岗位"的错配，让转型路径精确到天。' if is_zh else
        'CareerLens ist keine Jobbörse. Wir reparieren den Skill-Job-Mismatch und machen Transformationspfade präzise bis auf den Tag.',
        font_size=16, color=SLATE)

    # Four-page funnel
    pages = [
        ('1. 我是谁', 'Page 1: Wer bin ich?',
         '简历→AI画像→职业心电图→叠加态分析', 'CV→KI-Profil→Vital Signs→Superposition',
         '🔬'),
        ('2. 去哪里', 'Page 2: Wohin?',
         '9大行业×60+岗位→购物车选择目标', '9 Branchen × 60+ Rollen → Zielauswahl',
         '🗺'),
        ('3. 怎么去', 'Page 3: Wie?',
         '技能→工具→认证→GitHub工程挑战', 'Skills→Tools→Zertifikate→GitHub-Engineering-Challenge',
         '🎯'),
        ('4. 开始走', 'Page 4: Los!',
         '中德地图→城市→公司→招聘链接', 'CN/DE-Karte→Städte→Unternehmen→Stellenangebote',
         '🚀'),
    ]

    for i, (title_zh, title_de, desc_zh, desc_de, icon) in enumerate(pages):
        x = 1 + i * 2.9
        add_rect(slide, x, 2.5, 2.6, 4, WHITE)
        add_textbox(slide, x + 0.2, 2.7, 2.2, 0.5, icon, font_size=28)
        add_textbox(slide, x + 0.2, 3.3, 2.2, 0.5,
            title_zh if is_zh else title_de,
            font_size=16, color=BLUE, bold=True)
        add_textbox(slide, x + 0.2, 4, 2.2, 2,
            desc_zh if is_zh else desc_de,
            font_size=12, color=SLATE)

    # ==================== SLIDE 4: Core Innovation ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '核心创新："能力叠加态"分析引擎' if is_zh else
        'Kerninnovation: „Kompetenz-Superposition"-Analyse-Engine',
        font_size=32, color=DARK, bold=True)

    add_textbox(slide, 1, 1.5, 11, 0.8,
        'LinkedIn告诉你"你适合什么工作"。CareerLens告诉你"你离一个全新的身份只差什么"。' if is_zh else
        'LinkedIn sagt „Du passt zu diesem Job". CareerLens sagt „Du bist nur 1 Skill von einer neuen Identität entfernt".',
        font_size=18, color=BLUE, bold=True)

    # Example
    add_rect(slide, 1, 2.5, 11.3, 4.3, WHITE)
    tf = add_textbox(slide, 1.3, 2.7, 10.8, 0.5,
        '示例：一位被VW裁员的传统机械工程师' if is_zh else
        'Beispiel: Ein bei VW entlassener Maschinenbauingenieur',
        font_size=16, color=DARK, bold=True)
    add_para(tf,
        '现有技能：FMEA, SPC, CATIA V5, GD&T, VDA 6.3, Lean Manufacturing' if is_zh else
        'Vorhandene Skills: FMEA, SPC, CATIA V5, GD&T, VDA 6.3, Lean Manufacturing',
        font_size=13, color=SLATE, space_before=8)
    add_para(tf, '', font_size=8)
    add_para(tf,
        '⚛ 叠加态分析结果：' if is_zh else '⚛ Superposition-Analyseergebnis:',
        font_size=15, color=BLUE, bold=True, space_before=12)
    add_para(tf, '', font_size=4)

    paths = [
        ('方案A: +ISO 13485 → 医疗器械质量工程师',
         'Option A: +ISO 13485 → Medical Device Quality Engineer',
         'AI风险: 72%→18% | 薪资: +35% | 培训: 4周 | 增长: 🔥高',
         'KI-Risiko: 72%→18% | Gehalt: +35% | Training: 4 Wochen | Wachstum: 🔥Hoch'),
        ('方案B: +BMS+电池安全 → 新能源电池检测工程师',
         'Option B: +BMS+Batteriesicherheit → EV-Batterietestingenieur',
         'AI风险: 72%→15% | 薪资: +42% | 培训: 6周 | 增长: 🔥高',
         'KI-Risiko: 72%→15% | Gehalt: +42% | Training: 6 Wochen | Wachstum: 🔥Hoch'),
        ('方案C: +ROS2+视觉检测 → 机器人质量保证工程师',
         'Option C: +ROS2+Vision → Robotik-Qualitätsingenieur',
         'AI风险: 72%→8% | 薪资: +55% | 培训: 8周 | 增长: 🔥高',
         'KI-Risiko: 72%→8% | Gehalt: +55% | Training: 8 Wochen | Wachstum: 🔥Hoch'),
    ]
    for title_zh, title_de, detail_zh, detail_de in paths:
        add_para(tf, title_zh if is_zh else title_de,
            font_size=14, color=DARK, bold=True, space_before=10)
        add_para(tf, detail_zh if is_zh else detail_de,
            font_size=11, color=EMERALD, space_before=2)

    # ==================== SLIDE 5: Three-Phase Vision ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '三阶段战略：工具 → 通路 → 信用' if is_zh else
        'Drei-Phasen-Strategie: Tool → Kanal → Kredit',
        font_size=32, color=DARK, bold=True)

    phases = [
        ('Phase 1: 工具层', 'Phase 1: Tool',
         '2026 Q2-Q4', BLUE,
         ['免费技能分析+叠加态推荐', '开源核心代码获取信任', '积累用户数据和行业认知', '目标: 10,000注册用户'],
         ['Kostenlose Skill-Analyse + Superposition', 'Open Source Core für Vertrauen', 'Nutzerdaten + Branchenwissen sammeln', 'Ziel: 10.000 Registrierungen']),
        ('Phase 2: 通路层', 'Phase 2: Kanal',
         '2027 Q1-Q4', EMERALD,
         ['BA/KURSNET API对接→课程精准匹配', '培训机构佣金(¥750-2250/单)', 'Vector/TÜV认证推荐合作', '目标: BA试点合同 + 月收入€10K+'],
         ['BA/KURSNET-API → Kursmatching', 'Trainingsprovisionen (€200-600/Vermittlung)', 'Vector/TÜV Zertifizierungspartner', 'Ziel: BA-Pilotvertrag + €10K+/Monat']),
        ('Phase 3: 信用层', 'Phase 3: Kredit',
         '2028+', AMBER,
         ['GitHub工程挑战→工程信用认证', '企业直聘通道(绕过简历初筛)', '技能数据报告→卖给培训机构和政府', '目标: 成为"全球工程师技能资产交易所"'],
         ['GitHub-Engineering-Challenge → Zertifizierung', 'Direktvermittlung an Unternehmen', 'Skill-Datenreports für Behörden', 'Ziel: „Globale Engineer Skill Exchange"']),
    ]

    for i, (title_zh, title_de, timeline, color, items_zh, items_de) in enumerate(phases):
        x = 1 + i * 3.8
        add_rect(slide, x, 1.8, 3.5, 5, WHITE)
        add_textbox(slide, x + 0.2, 1.9, 3.1, 0.4,
            title_zh if is_zh else title_de,
            font_size=18, color=color, bold=True)
        add_textbox(slide, x + 0.2, 2.4, 3.1, 0.3,
            timeline, font_size=12, color=SLATE)

        tf = add_textbox(slide, x + 0.2, 2.9, 3.1, 3.5, '', font_size=12, color=SLATE)
        items = items_zh if is_zh else items_de
        for item in items:
            add_para(tf, f'• {item}', font_size=12, color=SLATE, space_before=8)

    # ==================== SLIDE 6: Business Model ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '盈利模式：三大收入来源' if is_zh else
        'Geschäftsmodell: Drei Einnahmequellen',
        font_size=32, color=DARK, bold=True)

    models = [
        ('🏛 B2G: 政府购买服务', '🏛 B2G: Regierungsaufträge',
         ['Arbeitsamt (BA) Bildungsgutschein', '每人评估费€50-200 或年度许可€100K-500K',
          '你的BA高管朋友 = 进入决策链', '德国300万失业汽车工人 = 巨大市场'],
         ['Arbeitsamt (BA) Bildungsgutschein-System', 'Assessment-Gebühr €50-200/Person oder Jahreslizenz €100K-500K',
          'Direkter Kontakt zu BA-Führungskraft', '300.000+ arbeitslose Automobilarbeiter = riesiger Markt'],
         '潜在年收入: €100K-500K', 'Potenzieller Jahresumsatz: €100K-500K'),
        ('🏢 B2B: 培训佣金+企业服务', '🏢 B2B: Trainingsprovisionen + Unternehmensservice',
         ['培训机构佣金: 10-20%/单 (¥750-2250)', 'Vector/TÜV认证推荐佣金',
          '企业直聘通道: 按入职收费', '行业技能报告: 卖给培训机构和猎头'],
         ['Trainingsprovisionen: 10-20% (€200-600/Vermittlung)', 'Vector/TÜV Zertifizierungsprovisionen',
          'Direktvermittlung: Erfolgsgebühr', 'Skill-Reports für Trainingsanbieter und Headhunter'],
         '潜在年收入: ¥50-100万', 'Potenzieller Jahresumsatz: €60K-120K'),
        ('👤 B2C: 增值服务', '👤 B2C: Premium-Services',
         ['AI深度职业报告: ¥99/€15', '专家代码审查: ¥299/€45',
          '1v1咨询: ¥299/€45', '金牌工程认证答辩: ¥999/€150'],
         ['KI-Karrierebericht: €15', 'Experten-Code-Review: €45',
          '1v1 Beratung: €45', 'Gold-Engineering-Zertifizierung: €150'],
         '潜在年收入: ¥30-80万', 'Potenzieller Jahresumsatz: €35K-100K'),
    ]

    for i, (title_zh, title_de, items_zh, items_de, rev_zh, rev_de) in enumerate(models):
        x = 1 + i * 3.8
        add_rect(slide, x, 1.6, 3.5, 4.8, WHITE)
        add_textbox(slide, x + 0.2, 1.7, 3.1, 0.5,
            title_zh if is_zh else title_de,
            font_size=16, color=BLUE, bold=True)

        tf = add_textbox(slide, x + 0.2, 2.3, 3.1, 3, '', font_size=12)
        items = items_zh if is_zh else items_de
        for item in items:
            add_para(tf, f'• {item}', font_size=11, color=SLATE, space_before=8)

        add_textbox(slide, x + 0.2, 5.6, 3.1, 0.4,
            rev_zh if is_zh else rev_de,
            font_size=13, color=EMERALD, bold=True)

    add_textbox(slide, 1, 6.6, 11, 0.5,
        '12个月综合盈利目标: €200K-700K（保守场景以BA一个试点合同为基础）' if is_zh else
        '12-Monats-Umsatzziel: €200K-700K (konservatives Szenario basierend auf einem BA-Pilotvertrag)',
        font_size=16, color=BLUE, bold=True)

    # ==================== SLIDE 7: Moat ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '护城河：为什么别人复制不了？' if is_zh else
        'Wettbewerbsvorteil: Warum ist das nicht kopierbar?',
        font_size=32, color=DARK, bold=True)

    moats = [
        ('⚛ 叠加态算法', '⚛ Superpositions-Algorithmus',
         '需要深厚的跨行业工程认知才能设计正确的技能桥梁映射。纯AI无法替代10年一线经验的行业判断力。' if is_zh else
         'Erfordert tiefes branchenübergreifendes Ingenieurwissen. KI allein kann 10 Jahre Branchenerfahrung nicht ersetzen.'),
        ('📊 专家验证数据飞轮', '📊 Experten-Daten-Flywheel',
         '用户越多→专家纠偏越多→路径越准→用户越多。这个数据资产随时间指数增长，后来者无法追上。' if is_zh else
         'Mehr Nutzer → mehr Experten-Feedback → präzisere Pfade → mehr Nutzer. Dieses Daten-Asset wächst exponentiell.'),
        ('🏛 政府关系', '🏛 Regierungsbeziehungen',
         '与BA高管的直接关系不是技术能复制的。政府采购有极高的信任门槛和合规要求。' if is_zh else
         'Direkter Kontakt zu BA-Führungskraft. Regierungsbeschaffung erfordert hohes Vertrauen und Compliance.'),
        ('🔬 工程信用体系', '🔬 Engineering-Credit-System',
         'GitHub工程挑战+隐藏测试+专家答辩=不可伪造的工程能力证明。LinkedIn做不到，因为它是社交平台，不是工程平台。' if is_zh else
         'GitHub-Challenge + Hidden Tests + Experten-Verteidigung = unfälschbarer Kompetenznachweis. LinkedIn kann das nicht – es ist Social, nicht Engineering.'),
        ('🌏 双市场定位', '🌏 Dual-Market-Positionierung',
         '同时覆盖中国和德国制造业市场。竞品通常只做一个市场。创始人中德双语+双市场认知。' if is_zh else
         'Gleichzeitig CN + DE Fertigungsmarkt. Wettbewerber bedienen nur einen. Gründer ist zweisprachig mit Marktkenntnis in beiden.'),
    ]

    for i, (title_zh, title_de, desc) in enumerate(moats):
        y = 1.6 + i * 1.1
        add_rect(slide, 1, y, 11.3, 0.95, WHITE)
        add_textbox(slide, 1.3, y + 0.05, 3, 0.4,
            title_zh if is_zh else title_de,
            font_size=14, color=BLUE, bold=True)
        add_textbox(slide, 1.3, y + 0.45, 10.8, 0.5,
            desc, font_size=11, color=SLATE)

    # ==================== SLIDE 8: Engineering Credit ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '终局愿景："全球工程师技能资产交易所"' if is_zh else
        'Endvision: „Global Engineer Skill Asset Exchange"',
        font_size=32, color=DARK, bold=True)

    add_textbox(slide, 1, 1.5, 11, 0.8,
        'Code is the only truth — 代码是唯一的真理' if is_zh else
        'Code is the only truth — der Code ist die einzige Wahrheit',
        font_size=20, color=BLUE, bold=True)

    # Three tiers
    tiers = [
        ('🥉', '铜牌: AI分析', '🥉 Bronze: KI-Analyse',
         '免费 · 即时', 'Kostenlos · Sofort',
         '简历→AI画像→叠加态推荐\n参考信息，不是证明', 'CV→KI-Profil→Superposition\nReferenzinfo, kein Nachweis'),
        ('🥈', '银牌: 项目完成', '🥈 Silber: Projekt abgeschlossen',
         '免费 · 需要时间', 'Kostenlos · Zeitinvestment',
         'Fork模板→完成代码→通过隐藏测试\n"这人真写过代码"', 'Fork Template→Code→Hidden Tests bestehen\n„Diese Person hat wirklich programmiert"'),
        ('🥇', '金牌: 专家验证', '🥇 Gold: Experten-Verifizierung',
         '付费 · 稀缺', 'Kostenpflichtig · Exklusiv',
         '银牌+30分钟口头答辩\n由行业专家亲自审核\n"不可伪造的工程能力证明"', 'Silber + 30min mündliche Verteidigung\nVon Branchenexperten geprüft\n„Unfälschbarer Engineering-Nachweis"'),
    ]

    for i, (icon, title_zh, title_de, sub_zh, sub_de, desc_zh, desc_de) in enumerate(tiers):
        x = 1 + i * 3.8
        bg_color = RGBColor(255, 251, 235) if i == 2 else WHITE
        add_rect(slide, x, 2.5, 3.5, 4.3, bg_color)
        add_textbox(slide, x + 0.2, 2.6, 3.1, 0.5,
            f'{title_zh}' if is_zh else f'{title_de}',
            font_size=18, color=DARK, bold=True)
        add_textbox(slide, x + 0.2, 3.2, 3.1, 0.3,
            sub_zh if is_zh else sub_de,
            font_size=12, color=AMBER if i == 2 else SLATE)
        add_textbox(slide, x + 0.2, 3.7, 3.1, 2.8,
            desc_zh if is_zh else desc_de,
            font_size=12, color=SLATE)

    # ==================== SLIDE 9: Market Size ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '市场规模' if is_zh else 'Marktgröße',
        font_size=32, color=DARK, bold=True)

    # TAM/SAM/SOM
    tf = add_textbox(slide, 1, 1.6, 11, 5, '', font_size=14)

    markets = [
        ('TAM (总可触达市场)', 'TAM (Total Addressable Market)',
         '全球工程师再培训市场 ~$50B/年' if is_zh else 'Globaler Ingenieur-Umschulungsmarkt ~$50 Mrd./Jahr',
         '德国Bildungsgutschein: €7B/年 | 中国职业培训补贴: ¥1000亿+/年' if is_zh else
         'DE Bildungsgutschein: €7 Mrd./Jahr | CN Berufsbildung: €14+ Mrd./Jahr'),
        ('SAM (可服务市场)', 'SAM (Serviceable Addressable Market)',
         '高端制造业工程师转型: ~$2B/年' if is_zh else 'Hightech-Fertigungs-Transformation: ~$2 Mrd./Jahr',
         '德国汽车+制造: 300万工人×平均€2000培训费 | 中国制造: 500万×¥5000' if is_zh else
         'DE Auto+Fertigung: 300K×€2.000 Durchschnitt | CN Fertigung: 5 Mio.×€700'),
        ('SOM (可获取市场)', 'SOM (Serviceable Obtainable Market)',
         '第一年目标: €200K-700K' if is_zh else 'Erstes Jahresziel: €200K-700K',
         '1个BA区域合同(€100-500K) + 培训佣金(€60-120K) + B2C(€35-100K)' if is_zh else
         '1 BA-Regionalvertrag (€100-500K) + Provisionen (€60-120K) + B2C (€35-100K)'),
    ]

    for title_zh, title_de, size_zh, size_de in markets:
        add_para(tf, title_zh if is_zh else title_de,
            font_size=18, color=BLUE, bold=True, space_before=16)
        add_para(tf, size_zh if is_zh else size_de,
            font_size=16, color=DARK, bold=True, space_before=4)
        add_para(tf, size_de.split('|')[0].strip() if is_zh else (size_de if '|' not in size_de else ''),
            font_size=12, color=SLATE, space_before=4)

    # ==================== SLIDE 10: Demo ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '产品演示：一个失业工程师的转型之路' if is_zh else
        'Produktdemo: Transformationspfad eines arbeitslosen Ingenieurs',
        font_size=32, color=DARK, bold=True)

    steps = [
        ('Step 1', '上传简历(PDF/文本)' if is_zh else 'CV hochladen (PDF/Text)',
         'AI自动识别: FMEA, SPC, CATIA, VDA 6.3等技能' if is_zh else
         'KI erkennt: FMEA, SPC, CATIA, VDA 6.3 etc.'),
        ('Step 2', '职业心电图即时显示' if is_zh else 'Karriere-Vitalzeichen sofort',
         'AI替代风险: 62% | 技能衰减: -15%/年 | 慕尼黑估价: €55-72K' if is_zh else
         'KI-Risiko: 62% | Skill-Decay: -15%/Jahr | München-Schätzung: €55-72K'),
        ('Step 3', '叠加态揭示"隐藏身份"' if is_zh else 'Superposition zeigt „versteckte Identität"',
         '"补1个技能(ISO 13485)→医疗器械质量工程师, 薪资+35%"' if is_zh else
         '„+1 Skill (ISO 13485) → Medical Device QE, Gehalt +35%"'),
        ('Step 4', '一键加入计划→行动路径' if is_zh else 'Ein-Klick-Plan → Aktionspfad',
         '技能清单+工具+认证+GitHub工程挑战→完成后获工程信用' if is_zh else
         'Skills + Tools + Zertifikate + GitHub-Challenge → Engineering Credit'),
        ('Step 5', '市场地图→直接投递' if is_zh else 'Marktkarte → Direkt bewerben',
         '慕尼黑西门子医疗正在招聘→带着工程信用直接面试' if is_zh else
         'Siemens Healthineers München sucht → mit Engineering Credit direkt bewerben'),
    ]

    for i, (step, title, desc) in enumerate(steps):
        y = 1.6 + i * 1.1
        add_rect(slide, 1, y, 11.3, 0.95, WHITE)
        add_textbox(slide, 1.3, y + 0.05, 1.2, 0.4, step, font_size=14, color=BLUE, bold=True)
        add_textbox(slide, 2.7, y + 0.05, 4, 0.4, title, font_size=14, color=DARK, bold=True)
        add_textbox(slide, 2.7, y + 0.45, 9.3, 0.4, desc, font_size=11, color=SLATE)

    # ==================== SLIDE 11: Roadmap ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_textbox(slide, 1, 0.5, 11, 0.8,
        '路线图' if is_zh else 'Roadmap',
        font_size=32, color=DARK, bold=True)

    roadmap = [
        ('2026 Q2', '已完成' if is_zh else 'Abgeschlossen', EMERALD,
         ['4页漏斗产品MVP', '叠加态分析引擎', '9行业×60+岗位数据', 'AI集成(通义千问)', '多模态简历输入(PDF/DOCX)'],
         ['4-Seiten-Funnel MVP', 'Superpositions-Engine', '9 Branchen × 60+ Rollen', 'KI-Integration (Qwen)', 'Multi-Input CV (PDF/DOCX)']),
        ('2026 Q3', '进行中' if is_zh else 'In Arbeit', BLUE,
         ['汽车行业数据深化(专家级)', 'GitHub工程挑战模板(3个)', 'BA Pitch + 试点提案', '在德华人社区推广'],
         ['Automotive-Daten vertiefen', 'GitHub-Challenge-Templates (3)', 'BA-Pitch + Pilotvorschlag', 'CN-Community in DE Marketing']),
        ('2026 Q4', '计划中' if is_zh else 'Geplant', AMBER,
         ['KURSNET API对接', 'Supabase后端(用户数据)', '专家纠偏系统上线', 'BA区域试点启动'],
         ['KURSNET-API-Integration', 'Supabase Backend', 'Experten-Feedback-System', 'BA-Regionalpilot Start']),
        ('2027', '愿景' if is_zh else 'Vision', SLATE,
         ['工程信用认证体系', '企业直聘通道', '行业技能数据报告', '扩展到IT/金融行业'],
         ['Engineering-Credit-System', 'Direktvermittlung', 'Branchen-Skill-Reports', 'Expansion: IT/Finance']),
    ]

    for i, (time, status_zh, color, items_zh, items_de) in enumerate(roadmap):
        x = 1 + i * 3
        add_rect(slide, x, 1.6, 2.7, 5, WHITE)
        add_textbox(slide, x + 0.15, 1.7, 2.4, 0.35, time, font_size=16, color=color, bold=True)
        add_textbox(slide, x + 0.15, 2.1, 2.4, 0.3,
            status_zh if is_zh else (status_zh),
            font_size=11, color=color)

        tf = add_textbox(slide, x + 0.15, 2.5, 2.4, 3.5, '', font_size=11)
        items = items_zh if is_zh else items_de
        for item in items:
            add_para(tf, f'• {item}', font_size=11, color=SLATE, space_before=6)

    # ==================== SLIDE 12: Team + Ask ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE)

    add_textbox(slide, 1, 0.8, 11, 0.8,
        '团队与资源' if is_zh else 'Team & Ressourcen',
        font_size=32, color=WHITE, bold=True)

    # Team
    add_rect(slide, 1, 2, 5.4, 3.5, RGBColor(30, 64, 175))
    tf = add_textbox(slide, 1.3, 2.1, 5, 0.5,
        '创始人: Michael Lin (林海青)' if is_zh else
        'Gründer: Michael Lin (林海青)',
        font_size=18, color=WHITE, bold=True)
    team_items = [
        'Magna 汽车电子质量管理 10年+' if is_zh else '10+ Jahre Automotive Quality bei Magna',
        'IATF 16949 主任审核员 / VDA 6.3' if is_zh else 'IATF 16949 Lead Auditor / VDA 6.3',
        '中德双语, 深耕德国制造业圈' if is_zh else 'Zweisprachig CN/DE, tief im DE-Fertigungssektor',
        '全栈开发: React/Next.js/TypeScript/AI' if is_zh else 'Full-Stack: React/Next.js/TypeScript/KI',
        'BA高管直接联系人' if is_zh else 'Direkter BA-Führungskontakt',
    ]
    for item in team_items:
        add_para(tf, f'• {item}', font_size=12, color=RGBColor(191, 219, 254), space_before=8)

    # Ask
    add_rect(slide, 6.8, 2, 5.4, 3.5, RGBColor(30, 64, 175))
    tf = add_textbox(slide, 7.1, 2.1, 5, 0.5,
        '我们需要什么' if is_zh else 'Was wir brauchen',
        font_size=18, color=WHITE, bold=True)
    ask_items = [
        ('BA试点合作: 1个区域3个月试点' if is_zh else 'BA-Pilotkooperation: 1 Region, 3 Monate'),
        ('培训机构合作: Vector / TÜV / AZAV机构' if is_zh else 'Trainingspartner: Vector / TÜV / AZAV-Träger'),
        ('行业专家: 3-5位资深工程师参与路径验证' if is_zh else 'Branchenexperten: 3-5 Senior Engineers für Pfadvalidierung'),
        ('在德华人社区渠道: 自媒体/微信群/行业协会' if is_zh else 'CN-Community-Kanäle: Social Media / WeChat / Fachverbände'),
    ]
    for item in ask_items:
        add_para(tf, f'• {item}', font_size=12, color=RGBColor(191, 219, 254), space_before=10)

    # Contact
    add_textbox(slide, 1, 6, 11, 1,
        'CareerLens — 让每个工程师看到自己的"隐藏身份"' if is_zh else
        'CareerLens — Zeige jedem Ingenieur seine „verborgene Identität"',
        font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 6.6, 11, 0.5,
        'careerlens.pages.dev · github.com/michaellinhq/careerlens · michaellinhq@gmail.com' ,
        font_size=14, color=RGBColor(148, 163, 184), alignment=PP_ALIGN.CENTER)

    return prs


if __name__ == '__main__':
    # Generate Chinese version
    prs_zh = build_deck('zh')
    prs_zh.save('CareerLens_商业计划书_中文版.pptx')
    print('✓ Chinese deck saved')

    # Generate German version
    prs_de = build_deck('de')
    prs_de.save('CareerLens_Geschaeftsplan_Deutsch.pptx')
    print('✓ German deck saved')
